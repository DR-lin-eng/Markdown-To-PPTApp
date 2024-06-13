import wx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os
import re

class MarkdownToPPTApp(wx.Frame):
    def __init__(self, parent, title):
        super(MarkdownToPPTApp, self).__init__(parent, title=title, size=(600, 400))
        
        panel = wx.Panel(self)
        vbox = wx.BoxSizer(wx.VERTICAL)

        self.bg_image_path = None

        self.bg_image_button = wx.Button(panel, label='选择底图')
        self.bg_image_button.Bind(wx.EVT_BUTTON, self.on_select_bg_image)
        vbox.Add(self.bg_image_button, flag=wx.EXPAND | wx.ALL, border=10)

        self.load_markdown_button = wx.Button(panel, label='导入Markdown文件')
        self.load_markdown_button.Bind(wx.EVT_BUTTON, self.on_load_markdown)
        vbox.Add(self.load_markdown_button, flag=wx.EXPAND | wx.ALL, border=10)

        self.text_ctrl = wx.TextCtrl(panel, style=wx.TE_MULTILINE)
        vbox.Add(self.text_ctrl, proportion=1, flag=wx.EXPAND | wx.ALL, border=10)

        self.convert_button = wx.Button(panel, label='转换为PPT')
        self.convert_button.Bind(wx.EVT_BUTTON, self.on_convert)
        vbox.Add(self.convert_button, flag=wx.EXPAND | wx.ALL, border=10)

        panel.SetSizer(vbox)

        self.Centre()
        self.Show(True)

    def on_select_bg_image(self, event):
        with wx.FileDialog(self, "选择底图文件", wildcard="图像文件 (*.jpg;*.jpeg;*.png)|*.jpg;*.jpeg;*.png",
                           style=wx.FD_OPEN | wx.FD_FILE_MUST_EXIST) as fileDialog:

            if fileDialog.ShowModal() == wx.ID_CANCEL:
                return

            self.bg_image_path = fileDialog.GetPath()
            wx.MessageBox(f'已选择底图: {self.bg_image_path}', '信息', wx.OK | wx.ICON_INFORMATION)

    def on_load_markdown(self, event):
        with wx.FileDialog(self, "导入Markdown文件", wildcard="Markdown文件 (*.md;*.txt)|*.md;*.txt",
                           style=wx.FD_OPEN | wx.FD_FILE_MUST_EXIST) as fileDialog:

            if fileDialog.ShowModal() == wx.ID_CANCEL:
                return

            markdown_path = fileDialog.GetPath()
            with open(markdown_path, 'r', encoding='utf-8') as file:
                markdown_content = file.read()
                self.text_ctrl.SetValue(markdown_content)
            wx.MessageBox(f'已导入文件: {markdown_path}', '信息', wx.OK | wx.ICON_INFORMATION)

    def on_convert(self, event):
        if not self.bg_image_path:
            wx.MessageBox('请先选择底图。', '错误', wx.OK | wx.ICON_ERROR)
            return

        markdown_content = self.text_ctrl.GetValue()
        if not markdown_content:
            wx.MessageBox('请输入Markdown文本。', '错误', wx.OK | wx.ICON_ERROR)
            return

        print("Markdown Content Received:")
        print(markdown_content)
        
        self.convert_markdown_to_ppt(markdown_content)
        wx.MessageBox('转换成功！', '信息', wx.OK | wx.ICON_INFORMATION)

    def convert_markdown_to_ppt(self, markdown_content):
        prs = Presentation()
        slides_data = markdown_content.split('---')[1:]  # Split by slide separators

        print("Slides Data:")
        print(slides_data)

        for slide_index, slide_data in enumerate(slides_data):
            print(f"\nProcessing Slide {slide_index + 1}")
            slide_layout = prs.slide_layouts[5]  # Using a blank slide layout
            slide = prs.slides.add_slide(slide_layout)
            self.set_background_image(slide, prs)
            self.handle_slide_data(slide, slide_data, prs.slide_width)

        prs.save('output_with_bg.pptx')
        print("Presentation saved as output_with_bg.pptx")

    def set_background_image(self, slide, prs):
        if self.bg_image_path:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
            slide.shapes.add_picture(self.bg_image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    def handle_slide_data(self, slide, slide_data, slide_width):
        lines = slide_data.strip().split('\n')
        title = None
        content = []

        print("Slide Data Lines:")
        for line in lines:
            print(line)
            if '###' in line:
                title = line.split('：')[1].strip()
            elif '- **标题**' in line:
                parts = line.split('**：')
                if len(parts) > 1:
                    title = parts[1].strip()
            elif '- **内容**' in line or line.startswith('    - '):
                if '- **内容**' in line:
                    parts = line.split('**：')
                    if len(parts) > 1:
                        content_line = parts[1].strip()
                        content.append(content_line)
                else:
                    content.append(line.strip())

        print(f"Title: {title}")
        print(f"Content: {content}")

        if title:
            self.add_title(slide, title, slide_width)
        if content:
            self.add_content(slide, content, slide_width)

    def add_title(self, slide, title, slide_width):
        left = Inches(0.5)
        top = Inches(0.5)
        width = slide_width - Inches(1)
        height = Inches(1.5)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        p = title_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.name = '华康方圆体W7'
        p.alignment = PP_ALIGN.CENTER
        title_frame.word_wrap = True  # 自动换行

    def add_content(self, slide, content_lines, slide_width):
        left = Inches(0.5)
        top = Inches(2)
        width = slide_width - Inches(1)
        height = Inches(5.0)
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        text_frame.word_wrap = True  # 自动换行
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        for line in content_lines:
            p = text_frame.add_paragraph()
            bold_parts = re.findall(r'\*\*(.*?)\*\*', line)
            plain_parts = re.split(r'\*\*(.*?)\*\*', line)
            for i, part in enumerate(plain_parts):
                run = p.add_run()
                run.text = part
                run.font.size = Pt(24)
                run.font.name = '华文新魏'
                if i % 2 == 1:
                    run.font.bold = True
            p.space_after = Pt(12)  # Add space after paragraphs for better separation
            p.alignment = PP_ALIGN.LEFT

if __name__ == '__main__':
    app = wx.App(False)
    frame = MarkdownToPPTApp(None, title='Markdown转PPT转换器')
    app.MainLoop()
