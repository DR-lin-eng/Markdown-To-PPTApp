import wx
import markdown2
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import re

class MarkdownToPPTApp(wx.Frame):
    def __init__(self, parent, title):
        super(MarkdownToPPTApp, self).__init__(parent, title=title, size=(400, 200))
        
        panel = wx.Panel(self)
        vbox = wx.BoxSizer(wx.VERTICAL)

        self.open_button = wx.Button(panel, label='Open Markdown/TXT File')
        self.open_button.Bind(wx.EVT_BUTTON, self.on_open_file)
        vbox.Add(self.open_button, flag=wx.EXPAND|wx.ALL, border=10)

        self.convert_button = wx.Button(panel, label='Convert to PPT')
        self.convert_button.Bind(wx.EVT_BUTTON, self.on_convert)
        vbox.Add(self.convert_button, flag=wx.EXPAND|wx.ALL, border=10)

        panel.SetSizer(vbox)

        self.markdown_file = None

        self.Centre()
        self.Show(True)
    
    def on_open_file(self, event):
        with wx.FileDialog(self, "Open Markdown or TXT file", wildcard="Markdown files (*.md)|*.md|Text files (*.txt)|*.txt",
                           style=wx.FD_OPEN | wx.FD_FILE_MUST_EXIST) as fileDialog:
            if fileDialog.ShowModal() == wx.ID_CANCEL:
                return

            self.markdown_file = fileDialog.GetPath()
            wx.MessageBox(f'Selected file: {self.markdown_file}', 'Info', wx.OK | wx.ICON_INFORMATION)
    
    def on_convert(self, event):
        if not self.markdown_file:
            wx.MessageBox('Please select a Markdown or TXT file first.', 'Error', wx.OK | wx.ICON_ERROR)
            return

        self.convert_markdown_to_ppt(self.markdown_file)
        wx.MessageBox('Conversion successful!', 'Info', wx.OK | wx.ICON_INFORMATION)
    
    def convert_markdown_to_ppt(self, markdown_file):
        with open(markdown_file, 'r', encoding='utf-8') as file:
            markdown_content = file.read()

        prs = Presentation()
        slides_data = markdown_content.split('---')[1:]  # Split by slide separators

        for slide_data in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using a blank slide layout
            title_text = None
            content_texts = []

            for line in slide_data.strip().split('\n'):
                if line.startswith('### 滑动'):
                    title_text = line.split('：')[1].strip()
                elif line.startswith('- **背景图片**'):
                    img_path = line.split('**：')[1].strip()
                    if os.path.exists(img_path):
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        pic = slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                elif line.startswith('- **标题**'):
                    title_text = line.split('**：')[1].strip()
                elif line.strip().startswith('- '):
                    content_texts.append(line.strip('- '))

            if title_text:
                title_shape = slide.shapes.title
                title_shape.text = title_text

            if content_texts:
                left = Inches(0.5)
                top = Inches(2)  # Adjust starting point for content below the title
                width = prs.slide_width - Inches(1)
                height = prs.slide_height - top - Inches(0.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame

                for text in content_texts:
                    p = tf.add_paragraph()
                    # Process and format bold text
                    bold_texts = re.findall(r'\*\*(.*?)\*\*', text)
                    for bold_text in bold_texts:
                        run = p.add_run()
                        run.text = bold_text
                        run.font.bold = True
                        text = text.replace(f"**{bold_text}**", "")
                    # Add remaining text
                    run = p.add_run()
                    run.text = text
                    p.font.size = Pt(18)
                    p.space_after = Pt(12)  # Add space after paragraphs for better separation

        output_file = markdown_file.replace('.md', '.pptx').replace('.txt', '.pptx')
        prs.save(output_file)

if __name__ == '__main__':
    app = wx.App(False)
    frame = MarkdownToPPTApp(None, title='Markdown to PPT Converter')
    app.MainLoop()
